"""
OCR & File Metadata Extraction — Unified Pipeline
===================================================
Handles text extraction + rich file metadata for ALL document types
flowing through any input node (upload, google_drive, dropbox, onedrive,
s3, ftp, url_scrape, email_inbox, webhook).

Capabilities:
  • PDF — direct text (PyMuPDF) + OCR fallback (tesseract) for scanned pages
  • Images (PNG/JPG/TIFF/BMP/GIF/WebP) — OCR via tesseract
  • DOCX/DOC — python-docx paragraph extraction
  • PPTX — python-pptx slide text extraction
  • XLSX/XLS — openpyxl cell extraction
  • HTML/XML/JSON/CSV/MD/TXT/RTF/ODT — format-specific parsers
  • File metadata: page count, word count, dimensions, language detection,
    text density, OCR confidence metrics

Usage:
    from clm.ocr_extraction import extract_all

    result = extract_all(file_obj, file_type='pdf')
    # result = {
    #     'direct_text': '...',
    #     'ocr_text': '...',
    #     'best_text': '...',
    #     'text_source': 'direct' | 'ocr' | 'none',
    #     'metadata': {
    #         'page_count': 5,
    #         'word_count': 2300,
    #         'char_count': 14500,
    #         'language': 'eng',
    #         'has_images': True,
    #         'has_tables': False,
    #         'ocr_confidence': 92.5,
    #         'text_density': 0.87,
    #         'dimensions': {'width': 612, 'height': 792, 'unit': 'pt'},
    #         'author': 'John Doe',
    #         'creation_date': '2024-03-15',
    #         'producer': 'Microsoft Word',
    #         'is_scanned': False,
    #         'extraction_method': 'pymupdf+ocr',
    #     }
    # }
"""

import io
import json
import logging
import os
import re
from datetime import datetime

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# File-type sets
# ---------------------------------------------------------------------------

IMAGE_EXTS = {'png', 'jpg', 'jpeg', 'gif', 'bmp', 'tiff', 'tif', 'webp'}
SPREADSHEET_EXTS = {'xlsx', 'xls'}
PRESENTATION_EXTS = {'pptx', 'ppt'}
TEXT_EXTS = {'txt', 'md', 'rtf', 'csv', 'json', 'xml', 'html', 'htm', 'svg', 'odt'}


# ---------------------------------------------------------------------------
# Main entry point
# ---------------------------------------------------------------------------

def extract_all(file_obj, file_type: str) -> dict:
    """
    Extract text + metadata from any file.
    Returns: {
        'direct_text': str,
        'ocr_text': str,
        'best_text': str,
        'text_source': 'direct' | 'ocr' | 'none',
        'metadata': dict,
    }
    """
    file_type = (file_type or '').lower().strip()
    direct_text = ''
    ocr_text = ''
    metadata = _base_metadata(file_type)

    # Ensure Django FieldFile objects are opened before reading
    opened = False
    if hasattr(file_obj, 'open') and hasattr(file_obj, 'field'):
        try:
            file_obj.open('rb')
            opened = True
        except Exception:
            pass

    try:
        if file_type == 'pdf':
            direct_text, ocr_text, metadata = _extract_pdf(file_obj, metadata)

        elif file_type in ('docx', 'doc'):
            direct_text, metadata = _extract_docx(file_obj, metadata)

        elif file_type in IMAGE_EXTS:
            ocr_text, metadata = _extract_image(file_obj, metadata)

        elif file_type in SPREADSHEET_EXTS:
            direct_text, metadata = _extract_spreadsheet(file_obj, file_type, metadata)

        elif file_type in PRESENTATION_EXTS:
            direct_text, metadata = _extract_presentation(file_obj, metadata)

        elif file_type == 'csv':
            direct_text, metadata = _extract_csv(file_obj, metadata)

        elif file_type == 'json':
            direct_text, metadata = _extract_json(file_obj, metadata)

        elif file_type in ('xml', 'svg'):
            direct_text, metadata = _extract_xml(file_obj, metadata)

        elif file_type in ('html', 'htm'):
            direct_text, metadata = _extract_html(file_obj, metadata)

        elif file_type == 'md':
            direct_text, metadata = _extract_markdown(file_obj, metadata)

        elif file_type == 'odt':
            direct_text, metadata = _extract_odt(file_obj, metadata)

        elif file_type in ('txt', 'rtf'):
            direct_text, metadata = _extract_text(file_obj, file_type, metadata)

        else:
            direct_text, metadata = _extract_fallback(file_obj, metadata)

    except Exception as e:
        logger.error(f"OCR extraction failed ({file_type}): {e}")
        metadata['extraction_error'] = str(e)
    finally:
        # Close Django FieldFile if we opened it
        if opened:
            try:
                file_obj.close()
            except Exception:
                pass

    # Decide best text
    text_source = 'none'
    best_text = ''
    if direct_text.strip() and len(direct_text.strip()) > 50:
        text_source = 'direct'
        best_text = direct_text
    elif ocr_text.strip() and len(ocr_text.strip()) > 50:
        text_source = 'ocr'
        best_text = ocr_text
    elif direct_text.strip():
        text_source = 'direct'
        best_text = direct_text
    elif ocr_text.strip():
        text_source = 'ocr'
        best_text = ocr_text

    # Common metadata from best text
    best = best_text.strip()
    if best:
        words = best.split()
        metadata['word_count'] = len(words)
        metadata['char_count'] = len(best)
        metadata['line_count'] = best.count('\n') + 1

        # Language detection (heuristic — top words)
        if not metadata.get('language'):
            metadata['language'] = _detect_language_heuristic(best)

    metadata['text_density'] = round(
        len(best) / max(metadata.get('page_count', 1), 1), 1
    ) if best else 0.0

    return {
        'direct_text': direct_text,
        'ocr_text': ocr_text,
        'best_text': best_text,
        'text_source': text_source,
        'metadata': metadata,
    }


def _base_metadata(file_type: str) -> dict:
    """Create base metadata dict with common defaults."""
    return {
        'file_format': file_type,
        'page_count': 0,
        'word_count': 0,
        'char_count': 0,
        'line_count': 0,
        'language': '',
        'has_images': False,
        'has_tables': False,
        'is_scanned': False,
        'ocr_confidence': 0.0,
        'text_density': 0.0,
        'extraction_method': '',
        'dimensions': None,
        'author': '',
        'creation_date': '',
        'producer': '',
    }


# ---------------------------------------------------------------------------
# PDF extraction — direct text + OCR for scanned pages
# ---------------------------------------------------------------------------

def _extract_pdf(file_obj, metadata: dict) -> tuple[str, str, dict]:
    """Extract from PDF: direct text + OCR for each page."""
    import fitz  # PyMuPDF

    file_obj.seek(0)
    pdf_bytes = file_obj.read()
    doc = fitz.open(stream=pdf_bytes, filetype='pdf')

    metadata['page_count'] = len(doc)
    metadata['extraction_method'] = 'pymupdf'

    # Document-level metadata
    pdf_meta = doc.metadata or {}
    metadata['author'] = pdf_meta.get('author', '') or ''
    metadata['producer'] = pdf_meta.get('producer', '') or ''
    metadata['creation_date'] = _parse_pdf_date(pdf_meta.get('creationDate', ''))

    if len(doc) > 0:
        page0 = doc[0]
        r = page0.rect
        metadata['dimensions'] = {
            'width': round(r.width, 1),
            'height': round(r.height, 1),
            'unit': 'pt',
        }

    # Extract direct text + track image-heavy pages for OCR
    direct_pages = []
    scanned_page_indices = []
    has_images = False

    for i, page in enumerate(doc):
        text = page.get_text().strip()
        direct_pages.append(text)

        # Check for images on this page
        image_list = page.get_images(full=True)
        if image_list:
            has_images = True

        # If a page has images but very little text, it's likely scanned
        if len(text) < 50 and image_list:
            scanned_page_indices.append(i)

    direct_text = '\n'.join(direct_pages)
    metadata['has_images'] = has_images

    # Check for tables (heuristic: look for tab-separated or pipe-separated patterns)
    if '|' in direct_text or '\t' in direct_text:
        metadata['has_tables'] = True

    # OCR for scanned/image pages
    ocr_text = ''
    ocr_confidences = []

    if scanned_page_indices or len(direct_text.strip()) < 50:
        # Either some pages are scanned, or the whole PDF has no text
        pages_to_ocr = scanned_page_indices if scanned_page_indices else range(len(doc))
        metadata['is_scanned'] = len(scanned_page_indices) == len(doc) or len(direct_text.strip()) < 50

        try:
            import pytesseract
            from PIL import Image

            ocr_pages = []
            for i in pages_to_ocr:
                page = doc[i]
                pix = page.get_pixmap(dpi=300)
                img = Image.open(io.BytesIO(pix.tobytes('png')))

                # Get OCR data with confidence
                ocr_data = pytesseract.image_to_data(img, output_type=pytesseract.Output.DICT)
                page_text = pytesseract.image_to_string(img)
                ocr_pages.append(page_text)

                # Compute page-level confidence
                confs = [int(c) for c in ocr_data['conf'] if str(c).lstrip('-').isdigit() and int(c) > 0]
                if confs:
                    ocr_confidences.extend(confs)

            ocr_text = '\n'.join(ocr_pages)
            metadata['extraction_method'] = 'pymupdf+ocr'
            logger.info(f"PDF OCR: {len(ocr_text.strip())} chars from {len(list(pages_to_ocr))} pages")

        except ImportError:
            logger.warning("OCR skipped: pytesseract not installed")
        except Exception as e:
            logger.warning(f"PDF OCR failed: {e}")

    if ocr_confidences:
        metadata['ocr_confidence'] = round(sum(ocr_confidences) / len(ocr_confidences), 1)

    doc.close()
    return direct_text, ocr_text, metadata


# ---------------------------------------------------------------------------
# Image extraction — OCR only
# ---------------------------------------------------------------------------

def _extract_image(file_obj, metadata: dict) -> tuple[str, dict]:
    """Extract text from image files via OCR."""
    metadata['page_count'] = 1
    metadata['has_images'] = True
    metadata['is_scanned'] = True
    metadata['extraction_method'] = 'ocr'

    try:
        from PIL import Image
        import pytesseract

        file_obj.seek(0)
        img = Image.open(file_obj)

        # Image dimensions
        metadata['dimensions'] = {
            'width': img.width,
            'height': img.height,
            'unit': 'px',
        }

        # Convert for OCR
        if img.mode not in ('L', 'RGB'):
            img = img.convert('RGB')

        # Get text
        text = pytesseract.image_to_string(img)

        # Get confidence data
        ocr_data = pytesseract.image_to_data(img, output_type=pytesseract.Output.DICT)
        confs = [int(c) for c in ocr_data['conf'] if str(c).lstrip('-').isdigit() and int(c) > 0]
        if confs:
            metadata['ocr_confidence'] = round(sum(confs) / len(confs), 1)

        # Detect orientation info
        try:
            osd = pytesseract.image_to_osd(img)
            rotation_match = re.search(r'Rotate: (\d+)', osd)
            if rotation_match:
                metadata['rotation'] = int(rotation_match.group(1))
            script_match = re.search(r'Script: (\w+)', osd)
            if script_match:
                metadata['script'] = script_match.group(1)
        except Exception:
            pass

        logger.info(f"Image OCR: {len(text.strip())} chars, confidence {metadata['ocr_confidence']}%")
        return text, metadata

    except ImportError:
        logger.warning("Image OCR skipped: pytesseract or Pillow not installed")
        return '', metadata
    except Exception as e:
        logger.warning(f"Image OCR failed: {e}")
        return '', metadata


# ---------------------------------------------------------------------------
# DOCX extraction
# ---------------------------------------------------------------------------

def _extract_docx(file_obj, metadata: dict) -> tuple[str, dict]:
    """Extract text from DOCX/DOC files."""
    metadata['extraction_method'] = 'python-docx'

    try:
        from docx import Document as DocxDocument

        file_obj.seek(0)
        doc = DocxDocument(file_obj)

        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        text = '\n'.join(paragraphs)

        # Count pages (approximate — DOCX doesn't store page count directly)
        # Estimate: ~3000 chars per page
        metadata['page_count'] = max(1, len(text) // 3000 + 1)

        # Check for tables
        if doc.tables:
            metadata['has_tables'] = True
            # Extract table text too
            table_texts = []
            for table in doc.tables:
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        table_texts.append(' | '.join(cells))
            if table_texts:
                text += '\n\n--- Tables ---\n' + '\n'.join(table_texts)

        # Check for images
        for rel in doc.part.rels.values():
            if 'image' in rel.reltype:
                metadata['has_images'] = True
                break

        # Core properties
        try:
            props = doc.core_properties
            metadata['author'] = props.author or ''
            if props.created:
                metadata['creation_date'] = props.created.strftime('%Y-%m-%d')
        except Exception:
            pass

        return text, metadata

    except ImportError:
        logger.warning("DOCX extraction skipped: python-docx not installed")
        return '', metadata
    except Exception as e:
        logger.warning(f"DOCX extraction failed: {e}")
        return '', metadata


# ---------------------------------------------------------------------------
# Presentation extraction (PPTX)
# ---------------------------------------------------------------------------

def _extract_presentation(file_obj, metadata: dict) -> tuple[str, dict]:
    """Extract text from PowerPoint files."""
    metadata['extraction_method'] = 'python-pptx'

    try:
        from pptx import Presentation

        file_obj.seek(0)
        prs = Presentation(file_obj)

        metadata['page_count'] = len(prs.slides)

        # Slide dimensions
        metadata['dimensions'] = {
            'width': round(prs.slide_width.pt, 1) if prs.slide_width else 0,
            'height': round(prs.slide_height.pt, 1) if prs.slide_height else 0,
            'unit': 'pt',
        }

        parts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_text = [f'--- Slide {i} ---']
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            slide_text.append(t)
                if shape.has_table:
                    metadata['has_tables'] = True
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            slide_text.append(' | '.join(cells))
                if shape.shape_type == 13:  # Picture
                    metadata['has_images'] = True
            parts.append('\n'.join(slide_text))

        text = '\n\n'.join(parts)
        return text, metadata

    except ImportError:
        logger.warning("PPTX extraction skipped: python-pptx not installed")
        return '', metadata
    except Exception as e:
        logger.warning(f"PPTX extraction failed: {e}")
        return '', metadata


# ---------------------------------------------------------------------------
# Spreadsheet extraction (XLSX/XLS)
# ---------------------------------------------------------------------------

def _extract_spreadsheet(file_obj, file_type: str, metadata: dict) -> tuple[str, dict]:
    """Extract text from Excel files."""
    metadata['extraction_method'] = 'openpyxl'
    metadata['has_tables'] = True

    try:
        import openpyxl

        file_obj.seek(0)
        wb = openpyxl.load_workbook(file_obj, read_only=True, data_only=True)

        metadata['page_count'] = len(wb.sheetnames)

        lines = []
        total_rows = 0
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            lines.append(f'--- Sheet: {sheet_name} ---')
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else '' for c in row]
                line = ' | '.join(cells).strip()
                if line and not all(c == '' for c in cells):
                    lines.append(line)
                    total_rows += 1

        wb.close()
        text = '\n'.join(lines)
        metadata['row_count'] = total_rows
        logger.info(f"Spreadsheet: {len(text.strip())} chars from {metadata['page_count']} sheets, {total_rows} rows")
        return text, metadata

    except ImportError:
        logger.warning("Spreadsheet extraction skipped: openpyxl not installed")
        return '', metadata
    except Exception as e:
        logger.warning(f"Spreadsheet extraction failed: {e}")
        return '', metadata


# ---------------------------------------------------------------------------
# CSV
# ---------------------------------------------------------------------------

def _extract_csv(file_obj, metadata: dict) -> tuple[str, dict]:
    """Extract text from CSV files."""
    import csv as csv_mod
    metadata['extraction_method'] = 'csv'
    metadata['has_tables'] = True
    metadata['page_count'] = 1

    try:
        file_obj.seek(0)
        raw = file_obj.read()
        text = raw.decode('utf-8', errors='replace') if isinstance(raw, bytes) else raw

        reader = csv_mod.reader(text.splitlines())
        lines = []
        for row in reader:
            lines.append(' | '.join(row))

        result = '\n'.join(lines)
        metadata['row_count'] = len(lines)
        return result, metadata

    except Exception as e:
        logger.warning(f"CSV extraction failed: {e}")
        return '', metadata


# ---------------------------------------------------------------------------
# JSON
# ---------------------------------------------------------------------------

def _extract_json(file_obj, metadata: dict) -> tuple[str, dict]:
    """Extract text from JSON files."""
    metadata['extraction_method'] = 'json'
    metadata['page_count'] = 1

    try:
        file_obj.seek(0)
        raw = file_obj.read()
        text = raw.decode('utf-8', errors='replace') if isinstance(raw, bytes) else raw
        data = json.loads(text)
        result = json.dumps(data, indent=2, ensure_ascii=False, default=str)
        return result, metadata

    except Exception as e:
        logger.warning(f"JSON extraction failed: {e}")
        return '', metadata


# ---------------------------------------------------------------------------
# XML / SVG
# ---------------------------------------------------------------------------

def _extract_xml(file_obj, metadata: dict) -> tuple[str, dict]:
    """Extract text from XML/SVG files."""
    metadata['extraction_method'] = 'xml'
    metadata['page_count'] = 1

    try:
        import xml.etree.ElementTree as ET

        file_obj.seek(0)
        raw = file_obj.read()
        text = raw.decode('utf-8', errors='replace') if isinstance(raw, bytes) else raw
        root = ET.fromstring(text)

        parts = []
        for elem in root.iter():
            if elem.text and elem.text.strip():
                parts.append(elem.text.strip())
            if elem.tail and elem.tail.strip():
                parts.append(elem.tail.strip())

        return '\n'.join(parts) if parts else text, metadata

    except Exception as e:
        logger.warning(f"XML extraction failed: {e}")
        try:
            file_obj.seek(0)
            raw = file_obj.read()
            return (raw.decode('utf-8', errors='replace') if isinstance(raw, bytes) else raw), metadata
        except Exception:
            return '', metadata


# ---------------------------------------------------------------------------
# HTML
# ---------------------------------------------------------------------------

def _extract_html(file_obj, metadata: dict) -> tuple[str, dict]:
    """Extract visible text from HTML."""
    metadata['extraction_method'] = 'html'
    metadata['page_count'] = 1

    try:
        file_obj.seek(0)
        raw = file_obj.read()
        text = raw.decode('utf-8', errors='replace') if isinstance(raw, bytes) else raw

        try:
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(text, 'html.parser')
            for tag in soup(['script', 'style', 'head']):
                tag.decompose()

            # Check for tables and images
            if soup.find('table'):
                metadata['has_tables'] = True
            if soup.find('img'):
                metadata['has_images'] = True

            result = soup.get_text(separator='\n', strip=True)
        except ImportError:
            result = re.sub(r'<[^>]+>', ' ', text)
            result = re.sub(r'\s+', ' ', result).strip()

        return result, metadata

    except Exception as e:
        logger.warning(f"HTML extraction failed: {e}")
        return '', metadata


# ---------------------------------------------------------------------------
# Markdown
# ---------------------------------------------------------------------------

def _extract_markdown(file_obj, metadata: dict) -> tuple[str, dict]:
    """Extract text from Markdown files."""
    metadata['extraction_method'] = 'markdown'
    metadata['page_count'] = 1

    try:
        file_obj.seek(0)
        raw = file_obj.read()
        text = raw.decode('utf-8', errors='replace') if isinstance(raw, bytes) else raw

        # Check for tables (markdown pipe tables)
        if re.search(r'\|.*\|.*\|', text):
            metadata['has_tables'] = True
        # Check for images
        if re.search(r'!\[.*\]\(.*\)', text):
            metadata['has_images'] = True

        # Strip markdown formatting but keep text
        text = re.sub(r'^#{1,6}\s+', '', text, flags=re.MULTILINE)
        text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
        text = re.sub(r'\*(.+?)\*', r'\1', text)
        text = re.sub(r'`(.+?)`', r'\1', text)
        text = re.sub(r'!\[.*?\]\(.*?\)', '', text)
        text = re.sub(r'\[(.+?)\]\(.*?\)', r'\1', text)

        return text, metadata

    except Exception as e:
        logger.warning(f"Markdown extraction failed: {e}")
        return '', metadata


# ---------------------------------------------------------------------------
# ODT
# ---------------------------------------------------------------------------

def _extract_odt(file_obj, metadata: dict) -> tuple[str, dict]:
    """Extract text from OpenDocument (.odt) files."""
    metadata['extraction_method'] = 'odt'

    try:
        import zipfile
        import xml.etree.ElementTree as ET

        file_obj.seek(0)
        with zipfile.ZipFile(file_obj) as zf:
            content_xml = zf.read('content.xml')

        root = ET.fromstring(content_xml)
        parts = []
        for p in root.iter('{urn:oasis:names:tc:opendocument:xmlns:text:1.0}p'):
            text = ''.join(p.itertext()).strip()
            if text:
                parts.append(text)

        result = '\n'.join(parts)
        metadata['page_count'] = max(1, len(result) // 3000 + 1)
        return result, metadata

    except Exception as e:
        logger.warning(f"ODT extraction failed: {e}")
        return '', metadata


# ---------------------------------------------------------------------------
# Plain text / RTF
# ---------------------------------------------------------------------------

def _extract_text(file_obj, file_type: str, metadata: dict) -> tuple[str, dict]:
    """Extract text from plain text or RTF."""
    metadata['extraction_method'] = file_type
    metadata['page_count'] = 1

    try:
        file_obj.seek(0)
        raw = file_obj.read()
        text = raw.decode('utf-8', errors='replace') if isinstance(raw, bytes) else raw

        if file_type == 'rtf':
            try:
                from striprtf.striprtf import rtf_to_text
                text = rtf_to_text(text)
                metadata['extraction_method'] = 'striprtf'
            except ImportError:
                logger.warning("RTF extraction: striprtf not installed, using raw text")

        return text, metadata

    except Exception as e:
        logger.warning(f"Text extraction failed ({file_type}): {e}")
        return '', metadata


# ---------------------------------------------------------------------------
# Fallback
# ---------------------------------------------------------------------------

def _extract_fallback(file_obj, metadata: dict) -> tuple[str, dict]:
    """Generic fallback — try to read as UTF-8."""
    metadata['extraction_method'] = 'fallback'
    metadata['page_count'] = 1

    try:
        file_obj.seek(0)
        raw = file_obj.read()
        if isinstance(raw, bytes):
            try:
                text = raw.decode('utf-8')
            except UnicodeDecodeError:
                text = raw.decode('latin-1', errors='replace')
        else:
            text = raw

        return text, metadata

    except Exception as e:
        logger.warning(f"Fallback extraction failed: {e}")
        return '', metadata


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _parse_pdf_date(date_str: str) -> str:
    """Parse PDF date format (D:20240315120000) to YYYY-MM-DD."""
    if not date_str:
        return ''
    # Strip the D: prefix and timezone info
    clean = re.sub(r"^D:", '', date_str).strip()
    clean = re.sub(r"[Z+']\d*'?\d*$", '', clean)
    try:
        if len(clean) >= 8:
            return f"{clean[:4]}-{clean[4:6]}-{clean[6:8]}"
    except Exception:
        pass
    return ''


def _detect_language_heuristic(text: str) -> str:
    """
    Simple language detection based on character frequency.
    Returns ISO 639-3 code: 'eng', 'spa', 'fra', 'deu', 'ita', 'por', etc.
    Falls back to 'eng' for Latin scripts.
    """
    if not text or len(text) < 20:
        return ''

    sample = text[:5000].lower()

    # Non-Latin scripts
    cjk = len(re.findall(r'[\u4e00-\u9fff]', sample))
    arabic = len(re.findall(r'[\u0600-\u06ff]', sample))
    devanagari = len(re.findall(r'[\u0900-\u097f]', sample))
    cyrillic = len(re.findall(r'[\u0400-\u04ff]', sample))
    korean = len(re.findall(r'[\uac00-\ud7af]', sample))
    japanese = len(re.findall(r'[\u3040-\u309f\u30a0-\u30ff]', sample))

    total_special = cjk + arabic + devanagari + cyrillic + korean + japanese
    if total_special > 50:
        if cjk > max(korean, japanese, arabic, devanagari, cyrillic):
            return 'zho'  # Chinese
        if korean > max(cjk, japanese):
            return 'kor'
        if japanese > cjk:
            return 'jpn'
        if arabic > max(cjk, devanagari, cyrillic):
            return 'ara'
        if devanagari > max(cjk, arabic, cyrillic):
            return 'hin'
        if cyrillic > max(cjk, arabic):
            return 'rus'

    # Latin script — word frequency heuristics
    words = sample.split()
    word_set = set(words[:500])

    # Spanish markers
    spanish = {'el', 'la', 'los', 'las', 'de', 'del', 'en', 'que', 'por', 'con', 'una', 'es', 'para'}
    if len(word_set & spanish) >= 5:
        return 'spa'

    # French markers
    french = {'le', 'la', 'les', 'de', 'des', 'du', 'un', 'une', 'est', 'dans', 'pour', 'avec', 'sur'}
    if len(word_set & french) >= 5:
        return 'fra'

    # German markers
    german = {'der', 'die', 'das', 'und', 'ist', 'von', 'ein', 'eine', 'mit', 'auf', 'den', 'für'}
    if len(word_set & german) >= 5:
        return 'deu'

    # Portuguese markers
    portuguese = {'de', 'da', 'do', 'dos', 'das', 'em', 'um', 'uma', 'para', 'com', 'por', 'não'}
    if len(word_set & portuguese) >= 5:
        return 'por'

    # Italian markers
    italian = {'il', 'la', 'di', 'che', 'non', 'una', 'per', 'del', 'della', 'sono', 'con'}
    if len(word_set & italian) >= 5:
        return 'ita'

    # Default: English
    return 'eng'
