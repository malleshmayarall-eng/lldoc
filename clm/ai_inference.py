"""
CLM AI Inference Module
========================
Singleton-managed NuExtract model for on-device metadata extraction.
Loads numind/NuExtract-1.5-tiny on MPS (Apple GPU), CUDA, or CPU.

Dual extraction:
  1. GLOBAL_CLM_TEMPLATE — standard legal fields every contract shares
  2. Workflow-specific template — built from rule-node field names

Both are extracted, standardized, and saved so rule-node filters work.

Usage:
    from clm.ai_inference import get_engine, extract_from_text

    # One-shot extraction from raw text
    result = extract_from_text("contract text...", {"field": ""})

    # Process a WorkflowDocument in-place (dual: global + workflow)
    from clm.ai_inference import extract_document
    result = extract_document(document, template)
"""
import io
import json
import logging
import os
import re
import threading
import time
from datetime import datetime
from decimal import Decimal, InvalidOperation
from typing import Any

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

MODEL_NAME = os.getenv('NUEXTRACT_MODEL', 'numind/NuExtract-1.5-tiny')
MAX_INPUT_CHARS = int(os.getenv('NUEXTRACT_MAX_INPUT_CHARS', '6000'))
MAX_NEW_TOKENS = int(os.getenv('NUEXTRACT_MAX_TOKENS', '2000'))
CHUNK_OVERLAP = int(os.getenv('NUEXTRACT_CHUNK_OVERLAP', '500'))
CONFIDENCE_THRESHOLD = float(os.getenv('NUEXTRACT_CONFIDENCE_THRESHOLD', '0.85'))


# ---------------------------------------------------------------------------
# Global CLM Template — standard legal/contract fields
# Split into logical groups for better extraction quality with small models.
# ---------------------------------------------------------------------------

GLOBAL_CLM_TEMPLATE = {
    "document_title": "",
    "document_type": "",
    "party_1_name": "",
    "party_1_role": "",
    "party_2_name": "",
    "party_2_role": "",
    "effective_date": "",
    "expiration_date": "",
    "execution_date": "",
    "contract_value": "",
    "currency": "",
    "payment_terms": "",
    "interest_rate": "",
    "governing_law": "",
    "jurisdiction": "",
    "termination_clause": "",
    "renewal_terms": "",
    "confidentiality": "",
    "indemnification": "",
    "dispute_resolution": "",
    "notice_address": "",
    "signatory_1": "",
    "signatory_2": "",
}

# Batched groups — small enough for the tiny model to handle well
_GLOBAL_TEMPLATE_BATCHES = [
    {   # Batch 1: Core identity (7 fields)
        "document_title": "",
        "document_type": "",
        "party_1_name": "",
        "party_1_role": "",
        "party_2_name": "",
        "party_2_role": "",
        "governing_law": "",
    },
    {   # Batch 2: Dates & money (7 fields)
        "effective_date": "",
        "expiration_date": "",
        "execution_date": "",
        "contract_value": "",
        "currency": "",
        "interest_rate": "",
        "payment_terms": "",
    },
    {   # Batch 3: Clauses & signatures (9 fields)
        "jurisdiction": "",
        "termination_clause": "",
        "renewal_terms": "",
        "confidentiality": "",
        "indemnification": "",
        "dispute_resolution": "",
        "notice_address": "",
        "signatory_1": "",
        "signatory_2": "",
    },
]


# ---------------------------------------------------------------------------
# Document Type Templates — type-specific extraction fields
# Each document type gets a tailored set of fields instead of the generic
# GLOBAL_CLM_TEMPLATE, producing much better extraction results.
# ---------------------------------------------------------------------------

DOCUMENT_TYPE_TEMPLATES = {
    'contract': {
        'label': 'Contract',
        'description': 'General contracts, service agreements, SLAs',
        'icon': '📄',
        'fields': {
            "document_title": "",
            "document_type": "",
            "party_1_name": "",
            "party_1_role": "",
            "party_2_name": "",
            "party_2_role": "",
            "effective_date": "",
            "expiration_date": "",
            "execution_date": "",
            "contract_value": "",
            "currency": "",
            "payment_terms": "",
            "governing_law": "",
            "jurisdiction": "",
            "termination_clause": "",
            "renewal_terms": "",
            "confidentiality": "",
            "indemnification": "",
            "dispute_resolution": "",
            "notice_address": "",
            "signatory_1": "",
            "signatory_2": "",
        },
    },
    'invoice': {
        'label': 'Invoice',
        'description': 'Invoices, bills, payment requests',
        'icon': '🧾',
        'fields': {
            "document_title": "",
            "invoice_number": "",
            "invoice_date": "",
            "due_date": "",
            "vendor_name": "",
            "vendor_address": "",
            "vendor_tax_id": "",
            "buyer_name": "",
            "buyer_address": "",
            "buyer_tax_id": "",
            "subtotal": "",
            "tax_amount": "",
            "tax_rate": "",
            "total_amount": "",
            "currency": "",
            "payment_terms": "",
            "payment_method": "",
            "bank_account": "",
            "purchase_order_number": "",
            "line_items_summary": "",
        },
    },
    'nda': {
        'label': 'NDA',
        'description': 'Non-disclosure / confidentiality agreements',
        'icon': '🔒',
        'fields': {
            "document_title": "",
            "agreement_type": "",
            "disclosing_party": "",
            "receiving_party": "",
            "effective_date": "",
            "expiration_date": "",
            "confidential_info_definition": "",
            "exclusions": "",
            "obligations": "",
            "term_duration": "",
            "governing_law": "",
            "jurisdiction": "",
            "return_of_materials": "",
            "remedies": "",
            "signatory_1": "",
            "signatory_2": "",
        },
    },
    'lease': {
        'label': 'Lease Agreement',
        'description': 'Property leases, rental agreements',
        'icon': '🏠',
        'fields': {
            "document_title": "",
            "landlord_name": "",
            "tenant_name": "",
            "property_address": "",
            "property_type": "",
            "lease_start_date": "",
            "lease_end_date": "",
            "monthly_rent": "",
            "security_deposit": "",
            "currency": "",
            "payment_due_day": "",
            "late_fee": "",
            "renewal_terms": "",
            "termination_notice_period": "",
            "permitted_use": "",
            "maintenance_responsibility": "",
            "governing_law": "",
            "signatory_1": "",
            "signatory_2": "",
        },
    },
    'employment': {
        'label': 'Employment Agreement',
        'description': 'Offer letters, employment contracts',
        'icon': '💼',
        'fields': {
            "document_title": "",
            "employer_name": "",
            "employee_name": "",
            "job_title": "",
            "department": "",
            "start_date": "",
            "employment_type": "",
            "base_salary": "",
            "currency": "",
            "pay_frequency": "",
            "bonus_structure": "",
            "benefits": "",
            "vacation_days": "",
            "probation_period": "",
            "termination_notice_period": "",
            "non_compete_clause": "",
            "confidentiality_clause": "",
            "governing_law": "",
            "signatory_1": "",
            "signatory_2": "",
        },
    },
    'purchase_order': {
        'label': 'Purchase Order',
        'description': 'POs, procurement orders',
        'icon': '🛒',
        'fields': {
            "document_title": "",
            "po_number": "",
            "order_date": "",
            "delivery_date": "",
            "buyer_name": "",
            "buyer_address": "",
            "supplier_name": "",
            "supplier_address": "",
            "line_items_summary": "",
            "subtotal": "",
            "tax_amount": "",
            "shipping_cost": "",
            "total_amount": "",
            "currency": "",
            "payment_terms": "",
            "shipping_method": "",
            "delivery_address": "",
            "special_instructions": "",
        },
    },
    'insurance': {
        'label': 'Insurance Policy',
        'description': 'Insurance policies, certificates of insurance',
        'icon': '🛡️',
        'fields': {
            "document_title": "",
            "policy_number": "",
            "insurer_name": "",
            "insured_name": "",
            "insured_address": "",
            "policy_type": "",
            "effective_date": "",
            "expiration_date": "",
            "coverage_amount": "",
            "deductible": "",
            "premium_amount": "",
            "payment_frequency": "",
            "currency": "",
            "beneficiary": "",
            "exclusions": "",
            "governing_law": "",
        },
    },
    'resume': {
        'label': 'Resume / CV',
        'description': 'Resumes, CVs, candidate profiles',
        'icon': '👤',
        'fields': {
            "candidate_name": "",
            "email": "",
            "phone": "",
            "location": "",
            "linkedin_url": "",
            "current_title": "",
            "current_company": "",
            "years_of_experience": "",
            "education_1": "",
            "education_2": "",
            "highest_degree": "",
            "university": "",
            "graduation_year": "",
            "technical_skills": "",
            "soft_skills": "",
            "programming_languages": "",
            "tools_and_frameworks": "",
            "certifications": "",
            "languages_spoken": "",
            "work_experience_1_company": "",
            "work_experience_1_title": "",
            "work_experience_1_duration": "",
            "work_experience_2_company": "",
            "work_experience_2_title": "",
            "work_experience_2_duration": "",
            "work_experience_3_company": "",
            "work_experience_3_title": "",
            "work_experience_3_duration": "",
            "summary": "",
        },
        'batches': [
            {  # Batch 1: Identity + contact
                "candidate_name": "",
                "email": "",
                "phone": "",
                "location": "",
                "linkedin_url": "",
                "current_title": "",
                "current_company": "",
            },
            {  # Batch 2: Education
                "years_of_experience": "",
                "education_1": "",
                "education_2": "",
                "highest_degree": "",
                "university": "",
                "graduation_year": "",
                "summary": "",
            },
            {  # Batch 3: Skills (most important — dedicated batch)
                "technical_skills": "",
                "soft_skills": "",
                "programming_languages": "",
                "tools_and_frameworks": "",
                "certifications": "",
                "languages_spoken": "",
            },
            {  # Batch 4: Work experience
                "work_experience_1_company": "",
                "work_experience_1_title": "",
                "work_experience_1_duration": "",
                "work_experience_2_company": "",
                "work_experience_2_title": "",
                "work_experience_2_duration": "",
                "work_experience_3_company": "",
                "work_experience_3_title": "",
                "work_experience_3_duration": "",
            },
        ],
    },
    'mou': {
        'label': 'Memorandum of Understanding',
        'description': 'MOUs, letters of intent',
        'icon': '🤝',
        'fields': {
            "document_title": "",
            "party_1_name": "",
            "party_2_name": "",
            "purpose": "",
            "effective_date": "",
            "duration": "",
            "responsibilities_party_1": "",
            "responsibilities_party_2": "",
            "financial_terms": "",
            "confidentiality": "",
            "termination_clause": "",
            "governing_law": "",
            "signatory_1": "",
            "signatory_2": "",
        },
    },
    'general': {
        'label': 'General Document',
        'description': 'Any document — uses standard legal fields',
        'icon': '📋',
        'fields': dict(GLOBAL_CLM_TEMPLATE),  # same as the original global template
    },
}


def get_template_for_type(document_type: str) -> dict:
    """
    Return the extraction template fields for a given document type.
    Falls back to GLOBAL_CLM_TEMPLATE if type is unknown.
    """
    entry = DOCUMENT_TYPE_TEMPLATES.get(document_type)
    if entry:
        return dict(entry['fields'])
    return dict(GLOBAL_CLM_TEMPLATE)


def get_batches_for_type(document_type: str) -> list[dict] | None:
    """
    Return custom batches for a document type, if defined.
    Types like 'resume' have hand-tuned batches that extract better
    than the automatic splitting.  Returns None if no custom batches.
    """
    entry = DOCUMENT_TYPE_TEMPLATES.get(document_type)
    if entry and 'batches' in entry:
        return list(entry['batches'])
    return None


def build_template_batches(template: dict, batch_size: int = 7) -> list[dict]:
    """
    Split a template into batches of *batch_size* fields each.
    Small batches work better with the tiny NuExtract model.
    """
    keys = list(template.keys())
    batches = []
    for i in range(0, len(keys), batch_size):
        batches.append({k: '' for k in keys[i:i + batch_size]})
    return batches


# ---------------------------------------------------------------------------
# Text extraction from files (PDF with OCR fallback / DOCX / TXT / Images / etc.)
# ---------------------------------------------------------------------------

# Image file extensions that should go through OCR
_IMAGE_EXTS = {'png', 'jpg', 'jpeg', 'gif', 'bmp', 'tiff', 'tif', 'webp'}

# Spreadsheet extensions
_SPREADSHEET_EXTS = {'xlsx', 'xls'}

# Plain-ish text formats we can read directly
_TEXT_EXTS = {'txt', 'md', 'rtf', 'csv', 'json', 'xml', 'html', 'svg', 'odt'}


def extract_text_from_file(file_obj, file_type: str) -> str:
    """
    Extract raw text from a file object based on type.
    Supports: PDF, DOCX/DOC, TXT, CSV, JSON, XML, HTML, Markdown,
    XLSX/XLS, images (OCR via tesseract), and more.
    Returns: the best available text string.
    """
    file_type = file_type.lower().strip()
    try:
        if file_type == 'pdf':
            direct, ocr = extract_pdf_dual(file_obj)
            return direct if direct.strip() and len(direct.strip()) > 50 else ocr

        elif file_type in ('docx', 'doc'):
            from docx import Document as DocxDocument
            file_obj.seek(0)
            doc = DocxDocument(file_obj)
            return '\n'.join(p.text for p in doc.paragraphs if p.text.strip())

        elif file_type in _IMAGE_EXTS:
            return _extract_text_from_image(file_obj)

        elif file_type in _SPREADSHEET_EXTS:
            return _extract_text_from_spreadsheet(file_obj, file_type)

        elif file_type == 'csv':
            return _extract_text_from_csv(file_obj)

        elif file_type == 'json':
            return _extract_text_from_json(file_obj)

        elif file_type in ('xml', 'svg'):
            return _extract_text_from_xml(file_obj)

        elif file_type == 'html':
            return _extract_text_from_html(file_obj)

        elif file_type == 'md':
            return _extract_text_from_markdown(file_obj)

        elif file_type == 'odt':
            return _extract_text_from_odt(file_obj)

        elif file_type in ('txt', 'rtf'):
            file_obj.seek(0)
            raw = file_obj.read()
            return raw.decode('utf-8', errors='replace') if isinstance(raw, bytes) else raw

        else:
            # Fallback: try to read as UTF-8 text
            file_obj.seek(0)
            raw = file_obj.read()
            text = raw.decode('utf-8', errors='replace') if isinstance(raw, bytes) else raw
            if text.strip():
                return text
            raise ValueError(f"Unsupported file type: {file_type}")
    except Exception as e:
        logger.error(f"Text extraction failed ({file_type}): {e}")
        return ''


def _extract_text_from_image(file_obj) -> str:
    """Extract text from image files using OCR (tesseract)."""
    try:
        import pytesseract
        from PIL import Image

        file_obj.seek(0)
        img = Image.open(file_obj)
        # Convert to RGB if needed (e.g. RGBA PNGs, palette images)
        if img.mode not in ('L', 'RGB'):
            img = img.convert('RGB')
        text = pytesseract.image_to_string(img)
        logger.info(f"Image OCR: {len(text.strip())} chars extracted")
        return text
    except ImportError:
        logger.warning("Image OCR skipped: pytesseract or Pillow not installed")
        return ''
    except Exception as e:
        logger.warning(f"Image OCR failed: {e}")
        return ''


def _extract_text_from_spreadsheet(file_obj, file_type: str) -> str:
    """Extract text from Excel files (.xlsx/.xls) — all sheets, all cells."""
    try:
        import openpyxl

        file_obj.seek(0)
        wb = openpyxl.load_workbook(file_obj, read_only=True, data_only=True)
        lines = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            lines.append(f"--- Sheet: {sheet_name} ---")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else '' for c in row]
                line = ' | '.join(cells).strip()
                if line and line != '|'.join([''] * len(cells)):
                    lines.append(line)
        wb.close()
        text = '\n'.join(lines)
        logger.info(f"Spreadsheet: {len(text.strip())} chars from {len(wb.sheetnames)} sheets")
        return text
    except ImportError:
        logger.warning("Spreadsheet extraction skipped: openpyxl not installed")
        # Fallback: try csv-like read
        return ''
    except Exception as e:
        logger.warning(f"Spreadsheet extraction failed: {e}")
        return ''


def _extract_text_from_csv(file_obj) -> str:
    """Extract text from CSV files."""
    import csv as csv_mod
    try:
        file_obj.seek(0)
        raw = file_obj.read()
        text = raw.decode('utf-8', errors='replace') if isinstance(raw, bytes) else raw
        # Parse and re-format for readability
        reader = csv_mod.reader(text.splitlines())
        lines = []
        for row in reader:
            lines.append(' | '.join(row))
        result = '\n'.join(lines)
        logger.info(f"CSV: {len(result.strip())} chars, {len(lines)} rows")
        return result
    except Exception as e:
        logger.warning(f"CSV extraction failed: {e}")
        return ''


def _extract_text_from_json(file_obj) -> str:
    """Extract text from JSON files — pretty-printed for AI readability."""
    try:
        file_obj.seek(0)
        raw = file_obj.read()
        text = raw.decode('utf-8', errors='replace') if isinstance(raw, bytes) else raw
        data = json.loads(text)
        result = json.dumps(data, indent=2, ensure_ascii=False, default=str)
        logger.info(f"JSON: {len(result.strip())} chars")
        return result
    except Exception as e:
        logger.warning(f"JSON extraction failed: {e}")
        return ''


def _extract_text_from_xml(file_obj) -> str:
    """Extract text content from XML/SVG files."""
    try:
        import xml.etree.ElementTree as ET

        file_obj.seek(0)
        raw = file_obj.read()
        text = raw.decode('utf-8', errors='replace') if isinstance(raw, bytes) else raw
        root = ET.fromstring(text)
        # Extract all text nodes
        parts = []
        for elem in root.iter():
            if elem.text and elem.text.strip():
                parts.append(elem.text.strip())
            if elem.tail and elem.tail.strip():
                parts.append(elem.tail.strip())
        result = '\n'.join(parts) if parts else text
        logger.info(f"XML: {len(result.strip())} chars")
        return result
    except Exception as e:
        logger.warning(f"XML extraction failed: {e}")
        # Return raw as fallback
        try:
            file_obj.seek(0)
            raw = file_obj.read()
            return raw.decode('utf-8', errors='replace') if isinstance(raw, bytes) else raw
        except Exception:
            return ''


def _extract_text_from_html(file_obj) -> str:
    """Extract visible text from HTML files."""
    try:
        file_obj.seek(0)
        raw = file_obj.read()
        text = raw.decode('utf-8', errors='replace') if isinstance(raw, bytes) else raw
        try:
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(text, 'html.parser')
            # Remove scripts and styles
            for tag in soup(['script', 'style', 'head']):
                tag.decompose()
            result = soup.get_text(separator='\n', strip=True)
        except ImportError:
            # Fallback: strip HTML tags with regex
            result = re.sub(r'<[^>]+>', ' ', text)
            result = re.sub(r'\s+', ' ', result).strip()
        logger.info(f"HTML: {len(result.strip())} chars")
        return result
    except Exception as e:
        logger.warning(f"HTML extraction failed: {e}")
        return ''


def _extract_text_from_markdown(file_obj) -> str:
    """Extract text from Markdown files (strip formatting)."""
    try:
        file_obj.seek(0)
        raw = file_obj.read()
        text = raw.decode('utf-8', errors='replace') if isinstance(raw, bytes) else raw
        # Strip common markdown syntax but keep the text
        text = re.sub(r'^#{1,6}\s+', '', text, flags=re.MULTILINE)  # headers
        text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)  # bold
        text = re.sub(r'\*(.+?)\*', r'\1', text)      # italic
        text = re.sub(r'`(.+?)`', r'\1', text)        # inline code
        text = re.sub(r'!\[.*?\]\(.*?\)', '', text)    # images
        text = re.sub(r'\[(.+?)\]\(.*?\)', r'\1', text)  # links
        logger.info(f"Markdown: {len(text.strip())} chars")
        return text
    except Exception as e:
        logger.warning(f"Markdown extraction failed: {e}")
        return ''


def _extract_text_from_odt(file_obj) -> str:
    """Extract text from OpenDocument (.odt) files."""
    try:
        import zipfile

        file_obj.seek(0)
        with zipfile.ZipFile(file_obj) as zf:
            content_xml = zf.read('content.xml')
        import xml.etree.ElementTree as ET
        root = ET.fromstring(content_xml)
        # Extract all text from text:p elements
        ns = {'text': 'urn:oasis:names:tc:opendocument:xmlns:text:1.0'}
        parts = []
        for p in root.iter('{urn:oasis:names:tc:opendocument:xmlns:text:1.0}p'):
            text = ''.join(p.itertext()).strip()
            if text:
                parts.append(text)
        result = '\n'.join(parts)
        logger.info(f"ODT: {len(result.strip())} chars")
        return result
    except Exception as e:
        logger.warning(f"ODT extraction failed: {e}")
        return ''


def extract_pdf_dual(file_obj) -> tuple[str, str]:
    """
    Extract BOTH direct text and OCR text from a PDF.
    Returns: (direct_text, ocr_text)
    Both are always attempted. Direct may be empty for scanned PDFs.
    """
    import fitz  # PyMuPDF

    file_obj.seek(0)
    pdf_bytes = file_obj.read()
    doc = fitz.open(stream=pdf_bytes, filetype='pdf')

    # Step 1: Direct text extraction
    direct_text = '\n'.join(page.get_text() for page in doc)
    logger.info(f"PDF direct text: {len(direct_text.strip())} chars")

    # Step 2: OCR (always attempt for scanned/image pages)
    ocr_text = ''
    try:
        import pytesseract
        from PIL import Image

        ocr_pages = []
        for page in doc:
            pix = page.get_pixmap(dpi=300)
            img = Image.open(io.BytesIO(pix.tobytes('png')))
            page_text = pytesseract.image_to_string(img)
            ocr_pages.append(page_text)

        ocr_text = '\n'.join(ocr_pages)
        logger.info(f"PDF OCR text: {len(ocr_text.strip())} chars from {len(ocr_pages)} pages")
    except ImportError:
        logger.warning("OCR skipped: pytesseract or Pillow not installed")
    except Exception as e:
        logger.warning(f"OCR failed: {e}")

    doc.close()
    return direct_text, ocr_text


def extract_dual_text_from_file(file_obj, file_type: str) -> tuple[str, str, str]:
    """
    Extract BOTH direct text and OCR text, plus determine which is best.
    Returns: (direct_text, ocr_text, text_source)
    text_source is 'direct', 'ocr', or 'none'.
    """
    file_type = file_type.lower().strip()
    direct_text = ''
    ocr_text = ''
    text_source = 'none'

    try:
        if file_type == 'pdf':
            direct_text, ocr_text = extract_pdf_dual(file_obj)
        elif file_type in ('docx', 'doc'):
            from docx import Document as DocxDocument
            file_obj.seek(0)
            doc = DocxDocument(file_obj)
            direct_text = '\n'.join(p.text for p in doc.paragraphs if p.text.strip())
        elif file_type == 'txt':
            file_obj.seek(0)
            raw = file_obj.read()
            direct_text = raw.decode('utf-8', errors='replace') if isinstance(raw, bytes) else raw
        elif file_type in ('jpg', 'jpeg', 'png', 'gif', 'bmp', 'tiff', 'webp'):
            # Images → OCR only
            try:
                import pytesseract
                from PIL import Image
                file_obj.seek(0)
                img = Image.open(file_obj)
                ocr_text = pytesseract.image_to_string(img)
            except ImportError:
                logger.warning("OCR skipped: pytesseract or Pillow not installed")
        elif file_type in ('xlsx', 'xls'):
            try:
                import openpyxl
                file_obj.seek(0)
                wb = openpyxl.load_workbook(file_obj, read_only=True, data_only=True)
                rows = []
                for ws in wb.worksheets:
                    rows.append(f"--- Sheet: {ws.title} ---")
                    for row in ws.iter_rows(values_only=True):
                        vals = [str(c) if c is not None else '' for c in row]
                        rows.append('\t'.join(vals))
                wb.close()
                direct_text = '\n'.join(rows)
            except ImportError:
                logger.warning("xlsx extraction skipped: openpyxl not installed")
        elif file_type == 'csv':
            import csv as csv_mod, io as _io
            file_obj.seek(0)
            raw = file_obj.read()
            text = raw.decode('utf-8', errors='replace') if isinstance(raw, bytes) else raw
            reader = csv_mod.reader(_io.StringIO(text))
            direct_text = '\n'.join('\t'.join(row) for row in reader)
        elif file_type in ('pptx', 'ppt'):
            try:
                from pptx import Presentation
                file_obj.seek(0)
                prs = Presentation(file_obj)
                parts = []
                for i, slide in enumerate(prs.slides, 1):
                    parts.append(f"--- Slide {i} ---")
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            parts.append(shape.text_frame.text)
                direct_text = '\n'.join(parts)
            except ImportError:
                logger.warning("pptx extraction skipped: python-pptx not installed")
        elif file_type in ('html', 'htm'):
            file_obj.seek(0)
            raw = file_obj.read()
            html_str = raw.decode('utf-8', errors='replace') if isinstance(raw, bytes) else raw
            try:
                from bs4 import BeautifulSoup
                direct_text = BeautifulSoup(html_str, 'html.parser').get_text(separator='\n')
            except ImportError:
                import re as _re
                direct_text = _re.sub(r'<[^>]+>', ' ', html_str)
        elif file_type in ('md', 'markdown'):
            file_obj.seek(0)
            raw = file_obj.read()
            direct_text = raw.decode('utf-8', errors='replace') if isinstance(raw, bytes) else raw
        elif file_type == 'rtf':
            try:
                from striprtf.striprtf import rtf_to_text
                file_obj.seek(0)
                raw = file_obj.read()
                rtf_str = raw.decode('utf-8', errors='replace') if isinstance(raw, bytes) else raw
                direct_text = rtf_to_text(rtf_str)
            except ImportError:
                logger.warning("RTF extraction skipped: striprtf not installed")
        else:
            # Generic fallback
            file_obj.seek(0)
            raw = file_obj.read()
            if isinstance(raw, bytes):
                try:
                    direct_text = raw.decode('utf-8')
                except UnicodeDecodeError:
                    direct_text = raw.decode('latin-1', errors='replace')
            else:
                direct_text = raw
    except Exception as e:
        logger.error(f"Text extraction failed ({file_type}): {e}")

    # Decide which to use: prefer direct if it has meaningful content
    if direct_text.strip() and len(direct_text.strip()) > 50:
        text_source = 'direct'
    elif ocr_text.strip() and len(ocr_text.strip()) > 50:
        text_source = 'ocr'
    else:
        text_source = 'none'

    return direct_text, ocr_text, text_source


# ---------------------------------------------------------------------------
# Prompt & chunking
# ---------------------------------------------------------------------------

def build_prompt(text: str, template: dict) -> str:
    """NuExtract v1.5 prompt format."""
    tmpl_json = json.dumps(template, indent=4)
    return f"<|input|>\n### Template:\n{tmpl_json}\n### Text:\n{text}\n\n<|output|>"


def chunk_text(text: str, max_length: int = MAX_INPUT_CHARS) -> list[str]:
    """Split long text into overlapping chunks."""
    if len(text) <= max_length:
        return [text]
    chunks, start = [], 0
    while start < len(text):
        chunks.append(text[start:start + max_length])
        start += max_length - CHUNK_OVERLAP
    return chunks


# ---------------------------------------------------------------------------
# JSON parser — robust against model quirks
# ---------------------------------------------------------------------------

def parse_json_output(text: str) -> dict:
    """Parse model output into a JSON dict, tolerating common quirks."""
    text = text.strip()
    # 1. Find the outermost { ... }
    brace_match = re.search(r'\{.*\}', text, re.DOTALL)
    if brace_match:
        try:
            return json.loads(brace_match.group())
        except json.JSONDecodeError:
            pass
    # 2. Whole text
    try:
        return json.loads(text)
    except json.JSONDecodeError:
        pass
    # 3. Fix trailing commas
    cleaned = re.sub(r',\s*}', '}', text)
    cleaned = re.sub(r',\s*]', ']', cleaned)
    try:
        return json.loads(cleaned)
    except json.JSONDecodeError:
        logger.warning(f"JSON parse failed: {text[:200]}")
        return {}


# ---------------------------------------------------------------------------
# Value Standardization — normalize extracted values so filters work
# ---------------------------------------------------------------------------

# Common date patterns the model might output
_DATE_PATTERNS = [
    (r'\d{4}-\d{2}-\d{2}', '%Y-%m-%d'),           # 2024-01-15
    (r'\d{2}/\d{2}/\d{4}', '%m/%d/%Y'),            # 01/15/2024
    (r'\d{2}-\d{2}-\d{4}', '%m-%d-%Y'),            # 01-15-2024
    (r'\d{1,2}\s+\w+\s+\d{4}', '%d %B %Y'),        # 15 January 2024
    (r'\w+\s+\d{1,2},?\s+\d{4}', '%B %d, %Y'),     # January 15, 2024
    (r'\w+\s+\d{1,2}\s+\d{4}', '%B %d %Y'),        # January 15 2024
]

# Fields that should be treated as dates
_DATE_FIELDS = {
    'effective_date', 'expiration_date', 'execution_date',
    'start_date', 'end_date', 'termination_date', 'renewal_date',
    'signing_date', 'contract_date', 'due_date', 'payment_date',
    'maturity_date', 'commencement_date',
}

# Fields that should be treated as monetary/numeric values
_NUMERIC_FIELDS = {
    'contract_value', 'total_value', 'payment_amount',
    'annual_value', 'monthly_payment', 'penalty_amount',
    'principal_amount', 'loan_amount', 'purchase_price',
    'interest_rate', 'rate', 'percentage',
}


def _standardize_date(value: str) -> str:
    """Try to parse a date string and return YYYY-MM-DD format."""
    if not value or not isinstance(value, str):
        return value
    value = value.strip()

    # Already standard
    if re.match(r'^\d{4}-\d{2}-\d{2}$', value):
        return value

    for pattern, fmt in _DATE_PATTERNS:
        m = re.search(pattern, value)
        if m:
            try:
                dt = datetime.strptime(m.group(), fmt)
                return dt.strftime('%Y-%m-%d')
            except ValueError:
                continue
    return value  # return as-is if no pattern matches


def _standardize_numeric(value: str) -> str:
    """
    Normalize monetary/numeric values to plain decimals.
    "$25,000.00" → "25000.00", "8%" → "8", "8% per annum" → "8"
    """
    if not value or not isinstance(value, str):
        return value
    v = value.strip()

    # Extract percentage: "8% per annum" → "8"
    pct_match = re.match(r'^([\d,.]+)\s*%', v)
    if pct_match:
        return pct_match.group(1).replace(',', '')

    # Strip currency symbols, commas, whitespace
    cleaned = re.sub(r'[^\d.\-]', '', v.replace(',', ''))
    if cleaned:
        try:
            d = Decimal(cleaned)
            return str(d)
        except InvalidOperation:
            pass
    return value


def _standardize_boolean(value: str) -> str:
    """Normalize yes/no/true/false to lowercase 'true'/'false'."""
    if not value or not isinstance(value, str):
        return value
    v = value.strip().lower()
    if v in ('yes', 'true', '1', 'y'):
        return 'true'
    if v in ('no', 'false', '0', 'n', 'none', 'n/a'):
        return 'false'
    return value


def standardize_extracted_data(data: dict) -> dict:
    """
    Post-process ALL extracted fields to standard formats so that
    rule-node condition evaluators (_eval_condition) work correctly.

    - Date fields → YYYY-MM-DD
    - Numeric/currency fields → plain decimal string
    - Everything else → stripped string
    """
    if not isinstance(data, dict):
        return data

    result = {}
    for field, value in data.items():
        if value is None or (isinstance(value, str) and not value.strip()):
            result[field] = value
            continue

        v = str(value).strip() if not isinstance(value, str) else value.strip()
        field_lower = field.lower()

        # Date fields
        if field_lower in _DATE_FIELDS or field_lower.endswith('_date'):
            result[field] = _standardize_date(v)
        # Numeric fields
        elif field_lower in _NUMERIC_FIELDS or field_lower.endswith('_value') or field_lower.endswith('_amount') or field_lower.endswith('_rate'):
            result[field] = _standardize_numeric(v)
        else:
            # General cleanup — trim, collapse whitespace
            result[field] = re.sub(r'\s+', ' ', v)

    return result


# ---------------------------------------------------------------------------
# Confidence heuristics (no-logit mode)
# ---------------------------------------------------------------------------

def estimate_confidence(template: dict, extracted: dict) -> dict[str, float]:
    """Heuristic per-field confidence when logits aren't available."""
    scores = {}
    for field in template:
        val = extracted.get(field)
        if not val or (isinstance(val, str) and not val.strip()):
            scores[field] = 0.0
            continue
        s = 0.7
        v = str(val)
        if len(v) > 3:
            s += 0.1
        if len(v) > 10:
            s += 0.05
        if field.endswith('_date') and re.match(r'\d{4}-\d{2}-\d{2}', v):
            s += 0.15
        elif field.endswith('_value') and re.match(r'[\d,.$]+', v):
            s += 0.15
        scores[field] = min(s, 1.0)
    return scores


# ---------------------------------------------------------------------------
# NuExtract Inference Engine (singleton)
# ---------------------------------------------------------------------------

class NuExtractEngine:
    """
    Thread-safe singleton that lazy-loads the NuExtract model once and
    serves all inference requests.
    """

    _instance = None
    _lock = threading.Lock()

    def __new__(cls):
        if cls._instance is None:
            with cls._lock:
                if cls._instance is None:
                    cls._instance = super().__new__(cls)
                    cls._instance._initialized = False
        return cls._instance

    def __init__(self):
        if self._initialized:
            return
        self._model = None
        self._tokenizer = None
        self._device = 'cpu'
        self._model_name = MODEL_NAME
        self._load_lock = threading.Lock()
        self._loaded = False
        self._load_error = None
        self._load_time = None
        self._inference_count = 0
        self._initialized = True

    # -- Status -------------------------------------------------------------

    @property
    def status(self) -> dict:
        return {
            'loaded': self._loaded,
            'model_name': self._model_name,
            'device': self._device,
            'load_error': str(self._load_error) if self._load_error else None,
            'load_time_seconds': self._load_time,
            'inference_count': self._inference_count,
        }

    # -- Load ---------------------------------------------------------------

    def _detect_device(self) -> str:
        import torch
        if torch.backends.mps.is_available():
            return 'mps'
        if torch.cuda.is_available():
            return 'cuda'
        return 'cpu'

    def load(self, force: bool = False):
        """Load the model. Thread-safe, idempotent unless force=True."""
        if self._loaded and not force:
            return
        with self._load_lock:
            if self._loaded and not force:
                return
            try:
                import torch
                from transformers import AutoModelForCausalLM, AutoTokenizer

                self._device = self._detect_device()
                logger.info(f"Loading {self._model_name} → {self._device}")
                t0 = time.time()

                self._tokenizer = AutoTokenizer.from_pretrained(
                    self._model_name, trust_remote_code=True,
                )
                dtype = torch.float16 if self._device in ('mps', 'cuda') else torch.float32
                self._model = AutoModelForCausalLM.from_pretrained(
                    self._model_name, trust_remote_code=True, torch_dtype=dtype,
                )
                self._model.to(self._device)
                self._model.eval()

                self._load_time = round(time.time() - t0, 2)
                self._loaded = True
                self._load_error = None
                logger.info(
                    f"Model loaded on {self._device} in {self._load_time}s "
                    f"({sum(p.numel() for p in self._model.parameters()):,} params)"
                )
            except Exception as e:
                self._load_error = e
                self._loaded = False
                logger.error(f"Model load failed: {e}")
                raise

    def ensure_loaded(self):
        """Ensure the model is loaded, loading it if needed."""
        if not self._loaded:
            self.load()

    # -- Inference ----------------------------------------------------------

    def _generate(self, prompt: str) -> str:
        """Tokenize, generate, decode — returns raw output text."""
        import torch
        self.ensure_loaded()

        inputs = self._tokenizer(
            prompt, return_tensors='pt', truncation=True, max_length=8192,
        )
        inputs = {k: v.to(self._device) for k, v in inputs.items()}

        with torch.no_grad():
            output_ids = self._model.generate(
                **inputs,
                max_new_tokens=MAX_NEW_TOKENS,
                do_sample=False,
            )

        prompt_len = inputs['input_ids'].shape[1]
        generated = output_ids[0][prompt_len:]
        self._inference_count += 1
        return self._tokenizer.decode(generated, skip_special_tokens=True)

    def extract(self, text: str, template: dict) -> dict:
        """
        Run NuExtract on text with the given template.
        Post-processes with value standardization so filters work.

        Returns:
            {
                "extracted_data": {...},       # standardized values
                "raw_extracted_data": {...},   # original model output
                "confidence": {...},
                "overall_confidence": float,
                "needs_review": bool,
                "chunks_processed": int,
            }
        """
        if not text or not text.strip():
            return {
                'extracted_data': {k: None for k in template},
                'raw_extracted_data': {k: None for k in template},
                'confidence': {k: 0.0 for k in template},
                'overall_confidence': 0.0,
                'needs_review': True,
                'chunks_processed': 0,
            }

        chunks = chunk_text(text)
        all_results: list[tuple[dict, dict[str, float]]] = []

        for chunk in chunks:
            prompt = build_prompt(chunk, template)
            raw = self._generate(prompt)
            parsed = parse_json_output(raw)
            conf = estimate_confidence(template, parsed)
            all_results.append((parsed, conf))

        # Merge: keep highest-confidence value per field
        merged_data = {}
        merged_conf = {}
        for field in template:
            best_val, best_c = None, 0.0
            for data, conf in all_results:
                c = conf.get(field, 0.0)
                v = data.get(field)
                if v and c > best_c:
                    best_val, best_c = v, c
            merged_data[field] = best_val
            merged_conf[field] = round(best_c, 4)

        conf_vals = [v for v in merged_conf.values() if v > 0]
        overall = round(sum(conf_vals) / len(conf_vals), 4) if conf_vals else 0.0

        # Standardize values for filter compatibility
        standardized = standardize_extracted_data(merged_data)

        return {
            'extracted_data': standardized,
            'raw_extracted_data': merged_data,
            'confidence': merged_conf,
            'overall_confidence': overall,
            'needs_review': any(v < CONFIDENCE_THRESHOLD for v in merged_conf.values()),
            'chunks_processed': len(chunks),
        }


# ---------------------------------------------------------------------------
# Module-level convenience API
# ---------------------------------------------------------------------------

_engine: NuExtractEngine | None = None


def get_engine() -> NuExtractEngine:
    """Get (or create) the singleton engine."""
    global _engine
    if _engine is None:
        _engine = NuExtractEngine()
    return _engine


def extract_from_text(text: str, template: dict) -> dict:
    """Extract metadata from raw text using the template."""
    engine = get_engine()
    return engine.extract(text, template)


def _merge_templates(global_template: dict, workflow_template: dict) -> dict:
    """
    Merge global CLM template with workflow-specific template.
    Workflow fields take priority (they may overlap with global).
    Returns: combined template with all unique fields.
    """
    merged = dict(global_template)
    merged.update(workflow_template)
    return merged


def _split_extraction_results(
    combined_data: dict,
    combined_confidence: dict,
    global_fields: set,
    workflow_fields: set,
) -> tuple[dict, dict, dict, dict]:
    """
    Split combined extraction results back into global and workflow-specific.
    Fields that appear in BOTH go to both dicts.
    Returns: (global_data, global_conf, workflow_data, workflow_conf)
    """
    global_data, global_conf = {}, {}
    workflow_data, workflow_conf = {}, {}

    for field, value in combined_data.items():
        conf = combined_confidence.get(field, 0.0)
        if field in global_fields:
            global_data[field] = value
            global_conf[field] = conf
        if field in workflow_fields:
            workflow_data[field] = value
            workflow_conf[field] = conf
        # Fields in both get saved to both
        if field not in global_fields and field not in workflow_fields:
            # Unexpected field from model output — put in global
            global_data[field] = value
            global_conf[field] = conf

    return global_data, global_conf, workflow_data, workflow_conf


# ---------------------------------------------------------------------------
# AI Field Discovery — Gemini analyses document text and chooses which
# metadata fields NuExtract should extract.  This replaces the static
# DOCUMENT_TYPE_TEMPLATES with a dynamic, document-aware field selection.
# ---------------------------------------------------------------------------

_FIELD_DISCOVERY_PROMPT = """You are an expert document analysis AI.
Given a snippet of text from a document, you must:
1. Identify the document type (one of: contract, invoice, nda, lease, employment, purchase_order, insurance, resume, mou, general)
2. Choose the best metadata fields to extract from this specific document.

You have access to these known document types and their standard fields:

DOCUMENT TYPES & STANDARD FIELDS:
- contract: document_title, document_type, party_1_name, party_1_role, party_2_name, party_2_role, effective_date, expiration_date, execution_date, contract_value, currency, payment_terms, governing_law, jurisdiction, termination_clause, renewal_terms, confidentiality, indemnification, dispute_resolution, notice_address, signatory_1, signatory_2
- invoice: document_title, invoice_number, invoice_date, due_date, vendor_name, vendor_address, vendor_tax_id, buyer_name, buyer_address, buyer_tax_id, subtotal, tax_amount, tax_rate, total_amount, currency, payment_terms, payment_method, bank_account, purchase_order_number, line_items_summary
- nda: document_title, agreement_type, disclosing_party, receiving_party, effective_date, expiration_date, confidential_info_definition, exclusions, obligations, term_duration, governing_law, jurisdiction, return_of_materials, remedies, signatory_1, signatory_2
- lease: document_title, landlord_name, tenant_name, property_address, property_type, lease_start_date, lease_end_date, monthly_rent, security_deposit, currency, payment_due_day, late_fee, renewal_terms, termination_notice_period, permitted_use, maintenance_responsibility, governing_law, signatory_1, signatory_2
- employment: document_title, employer_name, employee_name, job_title, department, start_date, employment_type, base_salary, currency, pay_frequency, bonus_structure, benefits, vacation_days, probation_period, termination_notice_period, non_compete_clause, confidentiality_clause, governing_law, signatory_1, signatory_2
- purchase_order: document_title, po_number, order_date, delivery_date, buyer_name, buyer_address, supplier_name, supplier_address, line_items_summary, subtotal, tax_amount, shipping_cost, total_amount, currency, payment_terms, shipping_method, delivery_address, special_instructions
- insurance: document_title, policy_number, insurer_name, insured_name, insured_address, policy_type, effective_date, expiration_date, coverage_amount, deductible, premium_amount, payment_frequency, currency, beneficiary, exclusions, governing_law
- resume: candidate_name, email, phone, location, linkedin_url, current_title, current_company, years_of_experience, education_1, education_2, highest_degree, university, graduation_year, technical_skills, soft_skills, programming_languages, tools_and_frameworks, certifications, languages_spoken, work_experience_1_company, work_experience_1_title, work_experience_1_duration, work_experience_2_company, work_experience_2_title, work_experience_2_duration, summary
- mou: document_title, party_1_name, party_2_name, purpose, effective_date, duration, responsibilities_party_1, responsibilities_party_2, financial_terms, confidentiality, termination_clause, governing_law, signatory_1, signatory_2

INSTRUCTIONS:
1. Analyse the text snippet carefully.
2. Determine which document type best matches.
3. Select the most relevant fields for THIS specific document — you can:
   - Use fields from the matching type template
   - REMOVE fields that clearly don't apply to this document
   - ADD custom fields (using snake_case naming) if the document contains
     important information not covered by the standard fields
4. Group fields into batches of 5-7 for optimal extraction quality.
   Each batch should contain related fields.

Return ONLY valid JSON (no markdown, no explanation):
{
  "document_type": "<detected type>",
  "confidence": <0.0-1.0>,
  "reasoning": "<one sentence explaining why>",
  "fields": {
    "<field_name>": "",
    ...
  },
  "batches": [
    {"<field1>": "", "<field2>": "", ...},
    {"<field3>": "", "<field4>": "", ...}
  ],
  "custom_fields_added": ["<any new fields not in standard templates>"]
}

RULES:
- Every field value must be an empty string ""
- Field names must be snake_case
- Include 10–30 fields depending on document complexity
- Batches should have 5–7 fields each
- Group related fields together in the same batch
- Always include the most important identifying fields first
- Naming convention matters for downstream type inference:
    * Fields ending in _date should contain dates
    * Fields ending in _amount, _value, _cost, _fee, _rent should contain numbers
    * Fields starting with is_ or has_ should be booleans
    * Fields named *_skills, *_items, *_certifications should be lists
"""


def discover_fields(text: str, document_type_hint: str = '') -> dict:
    """
    Use Gemini AI to analyse document text and determine the optimal
    metadata fields for extraction.

    Args:
        text: The extracted document text (best_text).
        document_type_hint: Optional hint from the workflow's input node.

    Returns:
        dict with keys: document_type, fields (template dict),
        batches (list of batch dicts), custom_fields_added, confidence.
        Falls back to static templates on any failure.
    """
    try:
        import google.generativeai as genai
    except ImportError:
        logger.warning("AI field discovery: google-generativeai not installed, using static templates")
        return _static_fallback(document_type_hint)

    from django.conf import settings as django_settings
    api_key = os.environ.get('GEMINI_API') or getattr(django_settings, 'GEMINI_API_KEY', '')
    if not api_key:
        logger.warning("AI field discovery: GEMINI_API not configured, using static templates")
        return _static_fallback(document_type_hint)

    # Use first 3000 chars — enough context for type detection + field selection
    snippet = text[:3000].strip()
    if not snippet:
        return _static_fallback(document_type_hint)

    genai.configure(api_key=api_key)

    model = genai.GenerativeModel(
        model_name='gemini-2.0-flash',
        system_instruction=_FIELD_DISCOVERY_PROMPT,
        generation_config=genai.GenerationConfig(
            temperature=0.2,  # Low temp for consistent field selection
            max_output_tokens=2048,
        ),
    )

    user_prompt = f"Analyse this document and suggest extraction fields:\n\n{snippet}"
    if document_type_hint:
        user_prompt += f"\n\nHint: the user expects this to be a '{document_type_hint}' document."

    try:
        response = model.generate_content(user_prompt)
        raw_text = response.text.strip()

        # Strip markdown fences
        if raw_text.startswith('```'):
            raw_text = re.sub(r'^```(?:json)?\s*', '', raw_text)
            raw_text = re.sub(r'\s*```$', '', raw_text)

        result = json.loads(raw_text)

        # Validate response structure
        if not isinstance(result, dict) or 'fields' not in result:
            logger.warning(f"AI field discovery returned unexpected structure, falling back")
            return _static_fallback(document_type_hint)

        # Ensure all field values are empty strings (NuExtract template format)
        fields = {k: '' for k, v in result['fields'].items() if isinstance(k, str)}
        batches = []
        for batch in result.get('batches', []):
            if isinstance(batch, dict):
                batches.append({k: '' for k in batch.keys()})

        # If no batches returned, auto-split
        if not batches:
            batches = build_template_batches(fields)

        detected_type = result.get('document_type', document_type_hint or 'general')

        logger.info(
            f"AI field discovery: detected type='{detected_type}', "
            f"{len(fields)} fields, {len(batches)} batches, "
            f"custom_fields={result.get('custom_fields_added', [])}"
        )

        return {
            'document_type': detected_type,
            'fields': fields,
            'batches': batches,
            'custom_fields_added': result.get('custom_fields_added', []),
            'confidence': result.get('confidence', 0.8),
            'reasoning': result.get('reasoning', ''),
            'source': 'ai_discovery',
        }

    except json.JSONDecodeError as e:
        logger.error(f"AI field discovery: invalid JSON response: {e}")
        return _static_fallback(document_type_hint)
    except Exception as e:
        logger.error(f"AI field discovery failed: {e}")
        return _static_fallback(document_type_hint)


def _static_fallback(document_type: str = '') -> dict:
    """
    Fallback to static templates when AI field discovery fails.
    Returns the same format as discover_fields().
    """
    template = get_template_for_type(document_type) if document_type else dict(GLOBAL_CLM_TEMPLATE)
    batches = get_batches_for_type(document_type) if document_type else None
    if not batches:
        batches = build_template_batches(template)

    return {
        'document_type': document_type or 'general',
        'fields': template,
        'batches': batches,
        'custom_fields_added': [],
        'confidence': 1.0,
        'reasoning': 'Using static template (AI discovery unavailable)',
        'source': 'static_template',
    }


# ---------------------------------------------------------------------------
# Gemini Re-extraction — fallback for fields NuExtract got wrong / low-conf
# ---------------------------------------------------------------------------

_GEMINI_EXTRACT_PROMPT = """You are a precise document metadata extractor.
Given a document's text and a list of metadata fields to extract, return
the values for each field based strictly on what the document contains.

OUTPUT: Return ONLY a valid JSON object — no markdown fences, no explanation,
no surrounding text. The JSON must be a flat object with the exact field names
as keys.

TYPE FORMATTING:
- STRING: Double-quoted JSON string with the exact value from the document.
- NUMBER: Bare numeric value without currency symbols, commas, or units.
  e.g. 50000 not "$50,000". For percentages, use the number: 5.5 not "5.5%".
- DATE: ISO 8601 format "YYYY-MM-DD". Convert any date format you find.
  e.g. "March 15, 2025" → "2025-03-15".
- BOOLEAN: Bare true or false (not quoted).
- LIST: JSON array of strings. e.g. ["Python", "React", "Node.js"].

RULES:
1. Extract ONLY values explicitly stated in the document.
2. Do NOT infer, assume, or fabricate any values.
3. If a field cannot be found, set its value to "" (empty string).
4. Prefer exact quotes from the document over paraphrasing.
5. For names and parties, use the full name as written in the document.
6. For monetary amounts, extract just the number with no formatting.
7. For date fields, always convert to YYYY-MM-DD regardless of source format.
"""


def gemini_extract_fields(text: str, fields: dict, document_type: str = '') -> dict:
    """
    Use Gemini AI to extract specific metadata fields from document text.
    Used as a fallback when NuExtract confidence is low on certain fields,
    or as the primary extractor for complex documents.

    Args:
        text: Document text (best_text).
        fields: Dict of {field_name: ""} — the template to extract.
        document_type: Optional type hint for better context.

    Returns:
        dict with keys: extracted_data, confidence (all 0.95 for Gemini).
        Returns empty results on failure.
    """
    if not fields:
        return {'extracted_data': {}, 'confidence': {}}

    try:
        import google.generativeai as genai
    except ImportError:
        logger.warning("Gemini extraction: google-generativeai not installed")
        return {'extracted_data': {}, 'confidence': {}}

    from django.conf import settings as django_settings
    api_key = os.environ.get('GEMINI_API') or getattr(django_settings, 'GEMINI_API_KEY', '')
    if not api_key:
        logger.warning("Gemini extraction: GEMINI_API not configured")
        return {'extracted_data': {}, 'confidence': {}}

    # Limit text to 4000 chars for cost/speed
    snippet = text[:4000].strip()
    if not snippet:
        return {'extracted_data': {}, 'confidence': {}}

    genai.configure(api_key=api_key)

    model = genai.GenerativeModel(
        model_name='gemini-2.0-flash',
        system_instruction=_GEMINI_EXTRACT_PROMPT,
        generation_config=genai.GenerationConfig(
            temperature=0.1,  # Very low for factual extraction
            max_output_tokens=2048,
        ),
    )

    field_list = ', '.join(fields.keys())
    user_prompt = (
        f"Extract these metadata fields from the document:\n"
        f"Fields: {field_list}\n"
    )
    if document_type:
        user_prompt += f"Document type: {document_type}\n"
    user_prompt += f"\n--- DOCUMENT TEXT ---\n{snippet}"

    try:
        response = model.generate_content(user_prompt)
        raw_text = response.text.strip()

        # Strip markdown fences
        if raw_text.startswith('```'):
            raw_text = re.sub(r'^```(?:json)?\s*', '', raw_text)
            raw_text = re.sub(r'\s*```$', '', raw_text)

        parsed = json.loads(raw_text)
        if not isinstance(parsed, dict):
            return {'extracted_data': {}, 'confidence': {}}

        # Only keep fields we asked for
        extracted = {}
        confidence = {}
        for field_name in fields:
            val = parsed.get(field_name, '')
            if val and str(val).strip():
                extracted[field_name] = str(val).strip()
                confidence[field_name] = 0.95  # Gemini is high confidence
            else:
                extracted[field_name] = ''
                confidence[field_name] = 0.0

        logger.info(
            f"Gemini extracted {sum(1 for v in extracted.values() if v)} / "
            f"{len(fields)} fields"
        )
        return {'extracted_data': extracted, 'confidence': confidence}

    except json.JSONDecodeError as e:
        logger.error(f"Gemini extraction: invalid JSON: {e}")
        return {'extracted_data': {}, 'confidence': {}}
    except Exception as e:
        logger.error(f"Gemini extraction failed: {e}")
        return {'extracted_data': {}, 'confidence': {}}


# ---------------------------------------------------------------------------
# AI Document Type Classification — Gemini detects document type from text
# ---------------------------------------------------------------------------

_VALID_DOCUMENT_TYPES = [
    'contract', 'invoice', 'nda', 'lease', 'employment',
    'purchase_order', 'insurance', 'resume', 'mou', 'general',
]

def classify_document_type(text: str) -> str:
    """
    Use Gemini to classify a document's type from its text content.
    Returns one of the valid document types, or 'general' if unsure.

    Used when the input node has no document_type configured, so the
    system can auto-detect and use the right extraction template.
    """
    if not text or not text.strip():
        return 'general'

    try:
        import google.generativeai as genai
    except ImportError:
        return 'general'

    from django.conf import settings as django_settings
    api_key = os.environ.get('GEMINI_API') or getattr(django_settings, 'GEMINI_API_KEY', '')
    if not api_key:
        return 'general'

    genai.configure(api_key=api_key)

    model = genai.GenerativeModel(
        model_name='gemini-2.0-flash',
        generation_config=genai.GenerationConfig(
            temperature=0.0,
            max_output_tokens=50,
        ),
    )

    snippet = text[:2000].strip()
    type_list = ', '.join(_VALID_DOCUMENT_TYPES)

    prompt = (
        f"Classify this document into exactly ONE of these types: {type_list}\n\n"
        f"Respond with ONLY the type name, nothing else.\n\n"
        f"--- DOCUMENT TEXT ---\n{snippet}"
    )

    try:
        response = model.generate_content(prompt)
        result = response.text.strip().lower().replace(' ', '_')
        # Validate it's one of the known types
        if result in _VALID_DOCUMENT_TYPES:
            logger.info(f"AI classified document as: {result}")
            return result
        # Fuzzy match (e.g. "purchase order" → "purchase_order")
        for dt in _VALID_DOCUMENT_TYPES:
            if dt in result or result in dt:
                logger.info(f"AI classified document as: {dt} (fuzzy: {result})")
                return dt
        logger.info(f"AI classification returned unknown type: {result}, defaulting to general")
        return 'general'
    except Exception as e:
        logger.error(f"AI document classification failed: {e}")
        return 'general'


def extract_document(document, workflow_template: dict, document_type: str = '', ai_discover: bool = False) -> dict:
    """
    Full dual-extraction pipeline for a WorkflowDocument:

    1. Extract text + OCR + file metadata via ocr_extraction.extract_all()
    1b. (Optional) AI Field Discovery — if ai_discover=True, call Gemini
        to analyse the document and choose optimal extraction fields.
        This replaces the static template with AI-selected fields.
    2. Run NuExtract PASS 1: type-specific global template
       (uses DOCUMENT_TYPE_TEMPLATES when *document_type* is set,
        otherwise falls back to GLOBAL_CLM_TEMPLATE)
       Uses custom batches for types like 'resume' that benefit from
       hand-tuned field grouping.
    3. Run NuExtract PASS 2: workflow_template (rule-node fields)
       — skipped if workflow_template is empty
    4. Post-process: list-aware parsing, value standardization
    5. Inject _text_snippet + _keywords into global_metadata so rule nodes
       can use 'contains' for full-text search / tag filtering.
    6. Save to document JSON fields + ExtractedField rows + ocr_metadata

    Two-pass approach: the tiny model (0.5B) handles smaller templates
    much better than one giant combined template.

    Returns dict with both global and workflow extraction results.
    """
    from .models import ExtractedField
    from .ocr_extraction import extract_all

    # 1. Unified text + OCR + metadata extraction
    ocr_result = extract_all(document.file, document.file_type)
    direct_text = ocr_result['direct_text']
    ocr_text = ocr_result['ocr_text']
    text_source = ocr_result['text_source']
    best_text = ocr_result['best_text']
    ocr_metadata = ocr_result['metadata']

    document.direct_text = direct_text
    document.ocr_text = ocr_text
    document.text_source = text_source
    document.original_text = best_text
    document.ocr_metadata = ocr_metadata
    document.extraction_status = 'processing'
    document.save(update_fields=[
        'direct_text', 'ocr_text', 'text_source',
        'original_text', 'ocr_metadata', 'extraction_status',
    ])

    if not best_text.strip():
        document.extraction_status = 'failed'
        document.save(update_fields=['extraction_status'])
        return {
            'global_metadata': {},
            'extracted_data': {},
            'confidence': {},
            'overall_confidence': 0.0,
            'text_source': text_source,
            'error': 'No text could be extracted from file.',
        }

    # 1a. Auto-detect document type if not set
    auto_classified = False
    if not document_type:
        document_type = classify_document_type(best_text)
        if document_type and document_type != 'general':
            auto_classified = True
            logger.info(f"Auto-classified doc {document.id} as '{document_type}'")
            # Store in global_metadata for downstream use
            gm = document.global_metadata or {}
            gm['_document_type'] = document_type
            gm['_auto_classified'] = True
            document.global_metadata = gm
            document.save(update_fields=['global_metadata'])

    # 1b. Determine extraction template — AI discovery or static
    discovery_info = None
    if ai_discover:
        # Let Gemini analyse the text and choose optimal fields + batches
        discovery_info = discover_fields(best_text, document_type_hint=document_type)
        global_template = discovery_info['fields']
        global_batches = discovery_info['batches']
        # Update document_type if AI detected something different
        if discovery_info['source'] == 'ai_discovery':
            document_type = discovery_info['document_type']
            logger.info(
                f"AI discovery: type={document_type}, {len(global_template)} fields, "
                f"{len(discovery_info.get('custom_fields_added', []))} custom fields"
            )
    else:
        # Static template selection
        global_template = get_template_for_type(document_type) if document_type else dict(GLOBAL_CLM_TEMPLATE)
        custom_batches = get_batches_for_type(document_type) if document_type else None
        global_batches = custom_batches or build_template_batches(global_template)

    engine = get_engine()

    # 2. PASS 1 — Global extraction in batches (smaller templates = better quality)
    global_data, global_raw, global_conf = {}, {}, {}
    for i, batch_template in enumerate(global_batches):
        logger.info(f"Pass 1.{i+1}: Global batch ({len(batch_template)} fields, type={document_type or 'general'})")
        try:
            batch_result = engine.extract(best_text, batch_template)
            global_data.update(batch_result['extracted_data'])
            global_raw.update(batch_result.get('raw_extracted_data', {}))
            global_conf.update(batch_result['confidence'])
        except Exception as e:
            logger.error(f"Global batch {i+1} failed for {document.id}: {e}")

    # 3. PASS 2 — Workflow-specific extraction (rule-node fields)
    workflow_data, workflow_raw, workflow_conf = {}, {}, {}
    workflow_fields = set(workflow_template.keys()) if workflow_template else set()

    if workflow_fields:
        # Remove fields already covered by the global template to avoid redundancy
        extra_fields = {k: v for k, v in workflow_template.items()
                        if k not in global_template}

        if extra_fields:
            logger.info(f"Pass 2: Workflow extraction ({len(extra_fields)} unique fields)")
            try:
                wf_result = engine.extract(best_text, extra_fields)
                workflow_data = wf_result['extracted_data']
                workflow_raw = wf_result.get('raw_extracted_data', {})
                workflow_conf = wf_result['confidence']
            except Exception as e:
                logger.error(f"Workflow extraction failed for {document.id}: {e}")

        # Copy fields from global that are also in workflow template
        for field in workflow_template:
            if field in global_template and field not in workflow_data:
                workflow_data[field] = global_data.get(field)
                workflow_conf[field] = global_conf.get(field, 0.0)
                workflow_raw[field] = global_raw.get(field)
    else:
        # No workflow template — copy global as workflow too
        workflow_data = dict(global_data)
        workflow_conf = dict(global_conf)
        workflow_raw = dict(global_raw)

    # 4. Post-process list-type fields (skills, certifications, etc.)
    #    NuExtract 0.5B sometimes returns messy comma-separated values or
    #    partial extractions. Normalise them so downstream consumers get
    #    clean, deduplicated, comma-separated lists.
    _LIST_FIELDS = {
        'technical_skills', 'soft_skills', 'programming_languages',
        'tools_and_frameworks', 'skills', 'certifications', 'languages',
        'keywords', 'tags', 'categories', 'signatories',
    }
    for data_dict in (global_data, workflow_data):
        for key in list(data_dict.keys()):
            if key.lower() in _LIST_FIELDS or key.lower().endswith('_skills'):
                val = data_dict[key]
                if isinstance(val, str) and val.strip():
                    # Split on commas, semicolons, pipes, or " and "
                    import re as _re
                    items = _re.split(r'[,;|]\s*|\s+and\s+', val)
                    items = [i.strip().strip('"\'') for i in items if i.strip()]
                    # Deduplicate while preserving order
                    seen = set()
                    unique = []
                    for item in items:
                        low = item.lower()
                        if low not in seen and len(low) > 1:
                            seen.add(low)
                            unique.append(item)
                    data_dict[key] = ', '.join(unique)

    # 4b. Gemini fallback — re-extract low-confidence or empty fields
    #     NuExtract 0.5B (tiny model) often produces empty or low-quality
    #     results. Use Gemini to fill in the blanks for better accuracy.
    _LOW_CONF = CONFIDENCE_THRESHOLD
    low_conf_fields = {}
    for field, conf in global_conf.items():
        val = global_data.get(field, '')
        if conf < _LOW_CONF or not str(val).strip():
            low_conf_fields[field] = ''
    for field, conf in workflow_conf.items():
        val = workflow_data.get(field, '')
        if conf < _LOW_CONF or not str(val).strip():
            low_conf_fields[field] = ''

    gemini_upgraded_fields = []
    if low_conf_fields:
        logger.info(
            f"Gemini fallback: {len(low_conf_fields)} low-confidence/empty fields "
            f"for doc {document.id}"
        )
        gemini_result = gemini_extract_fields(
            best_text, low_conf_fields, document_type=document_type,
        )
        gemini_data = gemini_result.get('extracted_data', {})
        gemini_conf = gemini_result.get('confidence', {})

        # Merge: only override if Gemini actually found a value
        for field, gemini_val in gemini_data.items():
            if not gemini_val:
                continue
            g_conf = gemini_conf.get(field, 0.0)
            # Override global fields
            if field in global_conf:
                old_val = global_data.get(field, '')
                old_conf = global_conf.get(field, 0.0)
                if g_conf > old_conf or not str(old_val).strip():
                    global_data[field] = gemini_val
                    global_conf[field] = g_conf
                    global_raw[field] = gemini_val
                    gemini_upgraded_fields.append(field)
            # Override workflow fields
            if field in workflow_conf:
                old_val = workflow_data.get(field, '')
                old_conf = workflow_conf.get(field, 0.0)
                if g_conf > old_conf or not str(old_val).strip():
                    workflow_data[field] = gemini_val
                    workflow_conf[field] = g_conf
                    workflow_raw[field] = gemini_val
                    if field not in gemini_upgraded_fields:
                        gemini_upgraded_fields.append(field)

        if gemini_upgraded_fields:
            logger.info(
                f"Gemini upgraded {len(gemini_upgraded_fields)} fields: "
                f"{gemini_upgraded_fields[:10]}"
            )

    # 5. Compute confidence
    global_conf_vals = [v for v in global_conf.values() if v > 0]
    workflow_conf_vals = [v for v in workflow_conf.values() if v > 0]
    global_overall = round(sum(global_conf_vals) / len(global_conf_vals), 4) if global_conf_vals else 0.0
    workflow_overall = round(sum(workflow_conf_vals) / len(workflow_conf_vals), 4) if workflow_conf_vals else 0.0

    all_conf = global_conf_vals + workflow_conf_vals
    overall = round(sum(all_conf) / len(all_conf), 4) if all_conf else 0.0

    # 6. Save to document
    # Preserve internal tracking keys (_source, source_hash, etc.) that were
    # set when the document was created.  AI extraction results must not
    # overwrite these — they are used for source filtering and deduplication.
    _preserved_keys = {'_source', 'source_hash'}
    existing_meta = document.global_metadata or {}
    preserved = {k: v for k, v in existing_meta.items() if k in _preserved_keys}
    merged_global = {**preserved, **global_data}

    # Inject text-as-metadata for rule-node 'contains' filtering.
    # _text_snippet: first 2000 chars of best_text for full-text search.
    # _keywords: auto-extracted top terms from the text for tag matching.
    _snippet = best_text[:2000].strip() if best_text else ''
    merged_global['_text_snippet'] = _snippet

    # Extract keywords: grab the most frequent meaningful words (4+ chars)
    import re as _re_kw
    from collections import Counter as _Counter
    _stop = {
        'this', 'that', 'with', 'from', 'have', 'been', 'will', 'were',
        'they', 'their', 'which', 'would', 'could', 'should', 'about',
        'there', 'shall', 'each', 'other', 'such', 'upon', 'into',
        'than', 'them', 'then', 'only', 'also', 'more', 'most', 'some',
        'what', 'when', 'your', 'does', 'made', 'make', 'just', 'very',
        'after', 'before', 'under', 'above', 'below', 'between', 'both',
        'being', 'those', 'these', 'here', 'where', 'while', 'same',
        'like', 'over', 'within', 'without', 'during', 'through',
        'any', 'all', 'and', 'the', 'for', 'are', 'but', 'not', 'you',
        'can', 'had', 'her', 'was', 'one', 'our', 'out', 'its', 'has',
        'may', 'per', 'who', 'how', 'nor',
    }
    _words = _re_kw.findall(r'\b[a-zA-Z]{4,}\b', best_text.lower())
    _freq = _Counter(w for w in _words if w not in _stop)
    _top_keywords = [w for w, _ in _freq.most_common(30)]
    merged_global['_keywords'] = ', '.join(_top_keywords)

    document.global_metadata = merged_global
    document.global_confidence = global_conf
    document.extracted_metadata = workflow_data
    document.extraction_confidence = workflow_conf
    document.overall_confidence = overall
    document.extraction_status = 'completed'
    document.save(update_fields=[
        'global_metadata', 'global_confidence',
        'extracted_metadata', 'extraction_confidence',
        'overall_confidence', 'extraction_status', 'updated_at',
    ])

    # 7. Save ExtractedField rows
    raw_data = {**global_raw, **workflow_raw}
    _save_extracted_fields(
        document=document,
        global_data=global_data,
        global_conf=global_conf,
        workflow_data=workflow_data,
        workflow_conf=workflow_conf,
        raw_data=raw_data,
    )

    result = {
        'global_metadata': global_data,
        'global_confidence': global_conf,
        'global_overall_confidence': global_overall,
        'extracted_data': workflow_data,
        'workflow_confidence': workflow_conf,
        'workflow_overall_confidence': workflow_overall,
        'overall_confidence': overall,
        'needs_review': any(v < CONFIDENCE_THRESHOLD for v in {**global_conf, **workflow_conf}.values()),
        'chunks_processed': 1,
        'text_source': text_source,
        'direct_text_length': len(direct_text.strip()),
        'ocr_text_length': len(ocr_text.strip()),
        'raw_extracted_data': raw_data,
        'ocr_metadata': ocr_metadata,
    }

    # Include AI discovery info if used
    if discovery_info:
        result['ai_discovery'] = {
            'detected_type': discovery_info.get('document_type', ''),
            'field_count': len(discovery_info.get('fields', {})),
            'custom_fields': discovery_info.get('custom_fields_added', []),
            'confidence': discovery_info.get('confidence', 0),
            'reasoning': discovery_info.get('reasoning', ''),
            'source': discovery_info.get('source', ''),
        }

    # Include Gemini fallback info
    if gemini_upgraded_fields:
        result['gemini_fallback'] = {
            'upgraded_fields': gemini_upgraded_fields,
            'upgraded_count': len(gemini_upgraded_fields),
        }

    # Include auto-classification info
    if auto_classified:
        result['auto_classification'] = {
            'detected_type': document_type,
            'auto_classified': True,
        }

    return result


def _save_extracted_fields(
    document,
    global_data: dict,
    global_conf: dict,
    workflow_data: dict,
    workflow_conf: dict,
    raw_data: dict,
):
    """
    Save/update individual ExtractedField rows for a document.
    Deletes old rows first, then bulk-creates new ones.
    """
    from .models import ExtractedField

    # Clear old rows for this document
    ExtractedField.objects.filter(document=document).delete()

    fields_to_create = []

    # Global fields
    for field_name, std_value in global_data.items():
        raw_val = raw_data.get(field_name, std_value) or ''
        conf = global_conf.get(field_name, 0.0)
        fields_to_create.append(ExtractedField(
            document=document,
            workflow=document.workflow,
            organization=document.organization,
            field_name=field_name,
            source='global',
            raw_value=str(raw_val) if raw_val else '',
            standardized_value=str(std_value) if std_value else '',
            display_value=str(raw_val) if raw_val else '',
            confidence=conf,
            needs_review=conf < CONFIDENCE_THRESHOLD,
        ))

    # Workflow-specific fields (skip duplicates already in global)
    for field_name, std_value in workflow_data.items():
        raw_val = raw_data.get(field_name, std_value) or ''
        conf = workflow_conf.get(field_name, 0.0)
        fields_to_create.append(ExtractedField(
            document=document,
            workflow=document.workflow,
            organization=document.organization,
            field_name=field_name,
            source='workflow',
            raw_value=str(raw_val) if raw_val else '',
            standardized_value=str(std_value) if std_value else '',
            display_value=str(raw_val) if raw_val else '',
            confidence=conf,
            needs_review=conf < CONFIDENCE_THRESHOLD,
        ))

    if fields_to_create:
        ExtractedField.objects.bulk_create(fields_to_create, ignore_conflicts=True)
        logger.info(f"Saved {len(fields_to_create)} ExtractedField rows for doc {document.id}")


def preload_model():
    """Pre-load the model (call at server start if desired)."""
    engine = get_engine()
    engine.load()
